"""PPTX document loader using python-pptx."""
from typing import Dict, Any, Optional, List
from pathlib import Path
import logging

logger = logging.getLogger(__name__)

# Check if python-pptx is available
PPTX_AVAILABLE = False
PPTX_UNAVAILABLE_REASON = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError as e:
    PPTX_UNAVAILABLE_REASON = f"python-pptx not installed: {str(e)}. Install with: pip install python-pptx"


class PPTXLoader:
    """Load PPTX documents using python-pptx."""
    
    def __init__(self):
        """Initialize PPTX loader."""
        self.name = "pptx"
    
    def extract_text(self, file_path: str) -> Dict[str, Any]:
        """
        Extract text from PPTX document.
        
        Args:
            file_path: Path to PowerPoint file
            
        Returns:
            Dictionary with extracted content
        """
        if not PPTX_AVAILABLE:
            return {
                "success": False,
                "loader": self.name,
                "error": PPTX_UNAVAILABLE_REASON
            }
        
        try:
            # Load presentation
            prs = Presentation(file_path)
            
            pages = []
            tables = []
            images = []
            text_parts = []
            table_index = 0
            image_index = 0
            
            for slide_idx, slide in enumerate(prs.slides):
                slide_text_parts = []
                slide_num = slide_idx + 1
                
                # Extract text from shapes
                for shape in slide.shapes:
                    # Text frames
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            text = paragraph.text.strip()
                            if text:
                                slide_text_parts.append(text)
                    
                    # Tables
                    if shape.has_table:
                        table = shape.table
                        headers = []
                        rows = []
                        
                        for row_idx, row in enumerate(table.rows):
                            row_data = [cell.text.strip() for cell in row.cells]
                            if row_idx == 0:
                                headers = row_data
                            else:
                                rows.append(row_data)
                        
                        tables.append({
                            "page_number": slide_num,
                            "table_index": table_index,
                            "headers": headers,
                            "rows": rows,
                            "caption": None,
                            "row_count": len(rows),
                            "col_count": len(headers)
                        })
                        table_index += 1
                    
                    # Images (just count them)
                    if hasattr(shape, 'image'):
                        images.append({
                            "page_number": slide_num,
                            "image_index": image_index,
                            "caption": None,
                            "alt_text": None
                        })
                        image_index += 1
                
                slide_text = "\n".join(slide_text_parts)
                text_parts.append(f"Slide {slide_num}:\n{slide_text}")
                
                pages.append({
                    "page_number": slide_num,
                    "text": slide_text,
                    "char_count": len(slide_text)
                })
            
            full_text = "\n\n".join(text_parts)
            
            # Get metadata
            core_props = prs.core_properties
            
            return {
                "success": True,
                "loader": self.name,
                "metadata": {
                    "title": core_props.title,
                    "author": core_props.author,
                    "format": "pptx",
                    "slide_count": len(prs.slides)
                },
                "pages": pages,
                "tables": tables,
                "images": images,
                "full_text": full_text,
                "total_pages": len(pages),
                "total_chars": len(full_text)
            }
            
        except Exception as e:
            logger.error(f"PPTX extraction failed: {str(e)}", exc_info=True)
            return {
                "success": False,
                "loader": self.name,
                "error": str(e)
            }
    
    def is_available(self) -> bool:
        """Check if loader is available."""
        return PPTX_AVAILABLE
    
    def get_unavailable_reason(self) -> Optional[str]:
        """Get reason why loader is unavailable."""
        return PPTX_UNAVAILABLE_REASON


# Global instance
pptx_loader = PPTXLoader()
